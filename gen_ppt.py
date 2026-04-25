from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────
BG     = RGBColor(0x0d, 0x11, 0x17)
CARD   = RGBColor(0x16, 0x1b, 0x22)
BORDER = RGBColor(0x30, 0x36, 0x3d)
TEXT   = RGBColor(0xe6, 0xed, 0xf3)
DIM    = RGBColor(0x8b, 0x94, 0x9e)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
PGBLUE = RGBColor(0x33, 0x67, 0x91)
TEAL   = RGBColor(0x06, 0xb6, 0xd4)
ABLUE  = RGBColor(0x01, 0x7c, 0xee)
BRONZE = RGBColor(0xcd, 0x9b, 0x1d)
SILVER = RGBColor(0xc0, 0xc0, 0xc0)
GOLD   = RGBColor(0xff, 0xd7, 0x00)
GREEN  = RGBColor(0x10, 0xb9, 0x81)
ORANGE = RGBColor(0xff, 0x66, 0x00)
RED    = RGBColor(0xff, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG
    return s

def tb(s, text, x, y, w, h, sz, bold=False, c=TEXT, al=PP_ALIGN.LEFT, it=False, font='Calibri'):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = c; r.font.italic = it; r.font.name = font
    return box

def code_tb(s, lines, x, y, w, h, c=RGBColor(0xa5,0xd6,0xff)):
    """Monospaced multi-line code block."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(13); r.font.name = 'Courier New'; r.font.color.rgb = c
    return box

def rect(s, x, y, w, h, fill=None, line=None, lw=1.5):
    sp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:    sp.fill.background()
    if line: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    else:    sp.line.fill.background()
    return sp

def hbar(s, color, y=0, h=0.07):
    rect(s, 0, y, 13.33, h, fill=color)

def header(s, title, subtitle, accent):
    hbar(s, accent)
    tb(s, title,    0.5, 0.14, 12.5, 0.7,  30, bold=True, c=TEXT)
    tb(s, subtitle, 0.5, 0.82, 12.5, 0.35, 15, c=DIM)
    rect(s, 0.5, 1.2, 12.33, 0.03, fill=accent)   # thin separator line

def card(s, x, y, w, h, color, title=None, title_sz=16):
    rect(s, x, y, w, h, fill=CARD, line=color, lw=2.0)
    if title:
        tb(s, title, x+0.2, y+0.12, w-0.4, 0.38, title_sz, bold=True, c=color)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
s1 = new_slide()
hbar(s1, PURPLE, y=0,    h=0.15)
hbar(s1, PGBLUE, y=7.35, h=0.15)

tb(s1, 'Real-Time Crypto Streaming', 0.8, 1.3,  11.73, 1.1, 46, bold=True, c=TEXT,  al=PP_ALIGN.CENTER)
tb(s1, '& Lakehouse Analytics',      0.8, 2.35, 11.73, 1.1, 46, bold=True, c=TEAL,  al=PP_ALIGN.CENTER)
tb(s1, 'Data Engineering Pipeline  ·  Project 4',
   0.8, 3.5, 11.73, 0.5, 20, c=DIM, al=PP_ALIGN.CENTER)
tb(s1, 'Coinbase WebSocket  ›  RabbitMQ  ›  PostgreSQL (Medallion)  ›  DuckLake  ›  Metabase',
   0.8, 4.05, 11.73, 0.45, 15, c=DIM, al=PP_ALIGN.CENTER)

badges = [('Python', PURPLE), ('RabbitMQ', ORANGE), ('PostgreSQL', PGBLUE),
          ('dbt', GREEN), ('DuckDB', GOLD), ('Metabase', TEAL), ('Airflow', ABLUE)]
bx = 1.3
for name, color in badges:
    bw = len(name) * 0.115 + 0.55
    rect(s1, bx, 4.9, bw, 0.5, fill=CARD, line=color, lw=1.5)
    tb(s1, name, bx+0.1, 4.93, bw-0.2, 0.44, 15, bold=True, c=color, al=PP_ALIGN.CENTER)
    bx += bw + 0.22


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Business Context
# ══════════════════════════════════════════════════════════════════════════
s2 = new_slide()
header(s2, 'Business Context', 'Why was this pipeline built?', PGBLUE)

# Client
card(s2, 0.4, 1.35, 5.9, 1.5, PGBLUE)
tb(s2, 'CLIENT', 0.6, 1.45, 5.5, 0.28, 11, bold=True, c=PGBLUE)
tb(s2, 'Boutique Quantitative Trading Firm', 0.6, 1.7, 5.5, 0.42, 18, bold=True, c=TEXT)
tb(s2, 'Need: Real-time VWAP & volatility dashboard', 0.6, 2.1, 5.5, 0.35, 14, c=DIM)

# Problem
card(s2, 0.4, 3.05, 5.9, 3.25, RED, title='THE PROBLEM')
for i, p in enumerate([
    '✗  Analysts query production DB directly',
    '✗  Heavy queries lock tables → missed trades',
    '✗  No separation of OLTP and OLAP workloads',
    '✗  System instability during market volatility',
]):
    tb(s2, p, 0.6, 3.62 + i*0.55, 5.5, 0.4, 15, c=TEXT)

# Solution
card(s2, 6.6, 1.35, 6.3, 5.95, GREEN, title='THE SOLUTION')
solutions = [
    ('WebSocket Ingestion',  'Coinbase → Python Producer',        PURPLE),
    ('Message Buffer',       'RabbitMQ absorbs velocity spikes',  ORANGE),
    ('Bronze Landing Zone',  'PostgreSQL stores raw JSON',        BRONZE),
    ('dbt Transformation',   'Silver (clean) → Gold (metrics)',   GREEN),
    ('Lakehouse Export',     'DuckDB: isolated analytical store', TEAL),
    ('BI Dashboard',         'Metabase on DuckLake data',         ABLUE),
]
for i, (title, desc, color) in enumerate(solutions):
    rect(s2, 6.8, 1.95+i*0.78, 5.9, 0.68, fill=BG, line=color, lw=1.2)
    tb(s2, title, 6.98, 2.0+i*0.78,  2.4, 0.3,  14, bold=True, c=color)
    tb(s2, desc,  9.42, 2.0+i*0.78,  3.2, 0.3,  13, c=DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════
s3 = new_slide()
header(s3, 'System Architecture', 'End-to-end data flow overview', PGBLUE)
s3.shapes.add_picture('architecture.png', Inches(0.3), Inches(1.28), Inches(12.73), Inches(6.0))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why RabbitMQ
# ══════════════════════════════════════════════════════════════════════════
s4 = new_slide()
header(s4, 'Architectural Reasoning: Why RabbitMQ?', 'Decoupling the WebSocket listener from the database', ORANGE)

# Without
card(s4, 0.4, 1.35, 5.9, 4.8, RED, title='WITHOUT RabbitMQ')
for i, (t, d) in enumerate([
    ('Tight Coupling',    'Producer blocked by slow DB writes'),
    ('Data Loss',         'DB crash = missed trades, no recovery'),
    ('Speed Mismatch',    'WebSocket rate > PostgreSQL write rate'),
    ('No Backpressure',   'Consumer overwhelmed at market open'),
]):
    tb(s4, f'✗  {t}', 0.6, 2.02+i*0.9, 5.5, 0.35, 16, bold=True, c=RED)
    tb(s4, f'     {d}', 0.6, 2.36+i*0.9, 5.5, 0.32, 13, c=DIM)

# With
card(s4, 7.0, 1.35, 5.9, 4.8, ORANGE, title='WITH RabbitMQ')
for i, (t, d) in enumerate([
    ('Shock Absorber',    'Queue buffers msgs during DB lag'),
    ('Fault Tolerance',   'Consumer crash → messages safely queued'),
    ('Full Decoupling',   'Producer/consumer scale independently'),
    ('Acknowledgement',   'Message removed only after INSERT'),
]):
    tb(s4, f'✓  {t}', 7.2, 2.02+i*0.9, 5.5, 0.35, 16, bold=True, c=ORANGE)
    tb(s4, f'     {d}', 7.2, 2.36+i*0.9, 5.5, 0.32, 13, c=DIM)

rect(s4, 0.4, 6.35, 12.5, 0.65, fill=CARD, line=ORANGE, lw=1.0)
tb(s4, 'Coinbase  ──►  Producer.py  ──►  [ RabbitMQ Queue ]  ──►  Consumer.py  ──►  PostgreSQL Bronze',
   0.6, 6.42, 12.1, 0.45, 16, bold=True, c=ORANGE, al=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Medallion Architecture
# ══════════════════════════════════════════════════════════════════════════
s5 = new_slide()
header(s5, 'Architectural Reasoning: Medallion Architecture', 'Three-layer transformation strategy with dbt', PGBLUE)

layers = [
    ('BRONZE', 'Raw Landing Zone', BRONZE, [
        'bronze.raw_trades',
        'Raw JSON — no transformations',
        'Complete audit trail',
        'Always replayable',
        '6 not_null tests on source',
    ]),
    ('SILVER', 'Cleaned & Trusted', SILVER, [
        'silver.trades',
        'DISTINCT ON(trade_id) deduplication',
        'CAST: NUMERIC price, TIMESTAMP time',
        'Filter: WHERE trade_id IS NOT NULL',
        '11 dbt tests: not_null + unique + accepted_values',
    ]),
    ('GOLD', 'Business Metrics', GOLD, [
        'gold.trade_metrics',
        'DATE_TRUNC("hour", time) aggregation',
        'VWAP, Volume, OHLC, buy_ratio',
        'FIRST/LAST_VALUE for price_open/close',
        '11 dbt tests on all 12 columns',
    ]),
]

for i, (name, tagline, color, bullets) in enumerate(layers):
    x = 0.4 + i * 4.28
    card(s5, x, 1.35, 4.05, 5.7, color)
    tb(s5, name,    x+0.15, 1.47, 3.75, 0.48, 24, bold=True, c=color, al=PP_ALIGN.CENTER)
    tb(s5, tagline, x+0.15, 1.93, 3.75, 0.3,  13, c=DIM,   al=PP_ALIGN.CENTER)
    rect(s5, x+0.15, 2.25, 3.75, 0.025, fill=color)
    tb(s5, bullets[0], x+0.2, 2.35, 3.65, 0.35, 14, bold=True, c=color)
    for j, b in enumerate(bullets[1:]):
        tb(s5, f'· {b}', x+0.2, 2.72+j*0.68, 3.65, 0.38, 13, c=TEXT)

tb(s5, 'dbt run  ▶', 4.5,  3.85, 0.78, 0.5, 14, bold=True, c=GREEN, al=PP_ALIGN.CENTER)
tb(s5, 'dbt run  ▶', 8.78, 3.85, 0.78, 0.5, 14, bold=True, c=GREEN, al=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Technical: JSON Handling
# ══════════════════════════════════════════════════════════════════════════
s6 = new_slide()
header(s6, 'Technical Deep Dive: JSON Payload Handling', 'Python handles ingest speed · dbt handles transformation quality', GREEN)

card(s6, 0.4, 1.35, 5.95, 5.7, PURPLE, title='Python  —  consumer.py')
tb(s6, 'Principle: ingest fast, transform later', 0.6, 1.82, 5.55, 0.28, 12, c=DIM, it=True)
rect(s6, 0.5, 2.12, 5.75, 4.68, fill=RGBColor(0x0a,0x0e,0x14), line=BORDER)
code_tb(s6, [
    'body = json.loads(msg_body)',
    '',
    'cursor.execute("""',
    '  INSERT INTO bronze.raw_trades',
    '    (trade_id, product_id,',
    '     price, size, side,',
    '     time, raw_data)',
    '  VALUES',
    '    (%s, %s, %s, %s, %s, %s, %s)',
    '""", (',
    '  body["trade_id"],',
    '  body["product_id"],',
    '  body["price"], ...))',
], 0.6, 2.22, 5.55, 4.5, c=RGBColor(0xa5,0xd6,0xff))

card(s6, 7.0, 1.35, 5.95, 5.7, GREEN, title='dbt  —  silver/trades.sql')
tb(s6, 'Principle: declarative, testable, version-controlled', 7.2, 1.82, 5.55, 0.28, 12, c=DIM, it=True)
rect(s6, 7.1, 2.12, 5.75, 4.68, fill=RGBColor(0x0a,0x0e,0x14), line=BORDER)
code_tb(s6, [
    'SELECT DISTINCT ON (trade_id)',
    '  trade_id,',
    '  product_id,',
    '  CAST(price AS NUMERIC) AS price,',
    '  CAST(size  AS NUMERIC) AS size,',
    '  side,',
    '  CAST(time AS TIMESTAMP) AS time,',
    '  raw_data',
    'FROM {{ source("bronze",',
    '               "raw_trades") }}',
    'WHERE trade_id IS NOT NULL',
    '  AND size IS NOT NULL',
    'ORDER BY trade_id, time',
], 7.22, 2.22, 5.55, 4.5, c=RGBColor(0x7e,0xe7,0x87))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — OLTP → OLAP
# ══════════════════════════════════════════════════════════════════════════
s7 = new_slide()
header(s7, 'Technical Deep Dive: OLTP → OLAP Migration', 'Separating transactional from analytical workloads', TEAL)

card(s7, 0.4, 1.35, 5.5, 4.75, PGBLUE, title='PostgreSQL  —  OLTP')
for i, (t, d) in enumerate([
    ('Row-store',      'Optimized for INSERT / UPDATE'),
    ('Transactional',  'ACID guarantees, row-level locking'),
    ('Write-optimized','Producer inserts 100+ rows / sec'),
    ('Our role',       'Landing zone + dbt transformations'),
]):
    tb(s7, t, 0.6, 2.02+i*0.82, 2.1, 0.32, 14, bold=True, c=PGBLUE)
    tb(s7, d, 2.75, 2.02+i*0.82, 2.9, 0.32, 13, c=DIM)

# Middle arrow / export description
rect(s7, 6.1, 1.35, 1.15, 4.75, fill=CARD, line=GOLD, lw=1.0)
for i, line in enumerate(['Gold', '──', 'pandas', '.read_sql()', '──', 'Parquet', '──', 'DuckLake', 'snapshot']):
    tb(s7, line, 6.13, 1.55+i*0.45, 1.0, 0.42, 12, bold=(line in ['Gold','DuckLake']),
       c=GOLD, al=PP_ALIGN.CENTER)

card(s7, 7.45, 1.35, 5.45, 4.75, TEAL, title='DuckDB  —  OLAP')
for i, (t, d) in enumerate([
    ('Columnar store',  '10-100x faster analytical scans'),
    ('Read-optimized',  'Zero write contention with producers'),
    ('Time-travel',     'Query any historical snapshot version'),
    ('Our role',        'Analyst-facing isolated lakehouse'),
]):
    tb(s7, t, 7.65, 2.02+i*0.82, 2.1, 0.32, 14, bold=True, c=TEAL)
    tb(s7, d, 9.8,  2.02+i*0.82, 2.9, 0.32, 13, c=DIM)

rect(s7, 0.4, 6.3, 12.5, 0.72, fill=CARD, line=GREEN, lw=1.0)
tb(s7, '✓  Key Insight: Analysts get 10-100x faster queries without ever touching the live trading system',
   0.6, 6.38, 12.1, 0.5, 15, bold=True, c=GREEN, al=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Live Demo
# ══════════════════════════════════════════════════════════════════════════
s8 = new_slide()
header(s8, 'Live System Demonstration', 'Start-to-finish pipeline walkthrough', ABLUE)

steps = [
    ('1', 'docker-compose up',        'Spin up 5 services: RabbitMQ :15672 · PostgreSQL :5433 · Airflow :8080 · Metabase :3000', ABLUE),
    ('2', 'python producer.py',       'Connect to Coinbase WebSocket → subscribe BTC/ETH/SOL-USD → push to RabbitMQ queue', PURPLE),
    ('3', 'RabbitMQ UI  →  :15672',   'Show crypto_trades queue filling in real time — messages accumulating', ORANGE),
    ('4', 'python consumer.py',       'Pull from RabbitMQ → INSERT into bronze.raw_trades → verify row count', PGBLUE),
    ('5', 'dbt run  +  dbt test',     'Build silver.trades (clean) → gold.trade_metrics (aggregated). All tests pass ✓', GREEN),
    ('6', 'python ducklake_export.py','Export Gold → DuckLake Parquet. Show time-travel snapshot history', GOLD),
    ('7', 'Metabase  →  :3000',       'Dashboard: VWAP by pair, hourly volume, buy/sell ratio live chart', TEAL),
    ('8', 'FAULT TOLERANCE DEMO',     'docker stop postgres → watch queue grow → docker start postgres → messages auto-consumed', RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    col, row = i % 2, i // 2
    x = 0.4 + col * 6.47
    y = 1.35 + row * 1.48
    rect(s8, x, y, 6.2, 1.32, fill=CARD, line=color, lw=1.5)
    tb(s8, num,   x+0.15, y+0.12, 0.48, 0.55, 26, bold=True, c=color, al=PP_ALIGN.CENTER)
    tb(s8, title, x+0.7,  y+0.12, 5.3,  0.38, 16, bold=True, c=color)
    tb(s8, desc,  x+0.7,  y+0.54, 5.3,  0.6,  12, c=DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Q&A
# ══════════════════════════════════════════════════════════════════════════
s9 = new_slide()
header(s9, 'Q&A  —  Anticipated Questions', 'Know your architectural decisions cold', PURPLE)

qas = [
    ('Why not WebSocket → PostgreSQL directly?',
     'No fault tolerance. If the DB is slow, the producer blocks and data is lost. RabbitMQ gives a durable buffer — messages persist even if the consumer crashes.',
     ORANGE),
    ('What happens when a dbt test fails?',
     'Pipeline stops. Data is NOT promoted to the next layer. Bronze always preserves the raw data — we can replay the transformation once the issue is fixed.',
     GREEN),
    ('Is DuckDB updated in real time?',
     'No — batch lakehouse with hourly refresh via Airflow. Tradeoff: data is max 1 hour old, but analysts get full read/write isolation from the production system.',
     TEAL),
    ('Why DuckDB instead of querying PostgreSQL Gold directly?',
     'Isolation: analysts cannot slow down the ingestion pipeline. DuckDB\'s columnar store is 10-100x faster for GROUP BY aggregations and time-series analysis.',
     PGBLUE),
]

for i, (q, a, color) in enumerate(qas):
    y = 1.35 + i * 1.48
    rect(s9, 0.4, y, 12.5, 1.32, fill=CARD, line=color, lw=1.5)
    tb(s9, f'Q:  {q}', 0.6, y+0.1,  12.1, 0.38, 15, bold=True, c=color)
    tb(s9, f'A:  {a}', 0.6, y+0.54, 12.1, 0.62, 13, c=TEXT)


# ── Save ──────────────────────────────────────────────────────────────────
prs.save('presentation.pptx')
print('Saved presentation.pptx  (9 slides)')
